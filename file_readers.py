from pathlib import Path
import typer
from rich.console import Console

import pdfplumber
from PyPDF2 import PdfReader
import docx
from pptx import Presentation

console = Console()


def _extract_pptx_text(file_path: str) -> str:
    """Extract text from a PowerPoint (.pptx), including shapes, tables, and speaker notes."""
    prs = Presentation(file_path)
    chunks = []

    for idx, slide in enumerate(prs.slides, start=1):
        # Slide title if present
        try:
            title = slide.shapes.title.text if slide.shapes.title else ""
            if title and title.strip():
                chunks.append(f"[Slide {idx} Title] {title.strip()}")
        except Exception:
            pass

        # Shapes (text boxes, placeholders)
        for shape in slide.shapes:
            # Plain text-bearing shapes
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    chunks.append(t)

            # Tables
            if getattr(shape, "has_table", False):
                try:
                    table = shape.table
                    for row in table.rows:
                        row_cells = []
                        for cell in row.cells:
                            cell_text = (cell.text or "").strip()
                            if cell_text:
                                row_cells.append(cell_text)
                        if row_cells:
                            chunks.append(" | ".join(row_cells))
                except Exception:
                    # Be permissive — skip malformed table entries
                    pass

        # Speaker notes (if present)
        try:
            if slide.has_notes_slide:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes_text:
                    chunks.append(f"[Slide {idx} Notes] {notes_text}")
        except Exception:
            pass

    return "\n".join(chunks)


def load_text(file_path: str) -> str:
    """Load text content from various file formats."""
    path_obj = Path(file_path)
    ext = path_obj.suffix.lower()

    if ext == ".txt":
        return path_obj.read_text(encoding="utf-8")

    elif ext == ".pdf":
        try:
            with pdfplumber.open(str(path_obj)) as pdf:
                return "\n".join(p.extract_text() or "" for p in pdf.pages)
        except Exception as e:
            console.print(f"[yellow]pdfplumber failed: {e}, trying PyPDF2...[/yellow]")

        try:
            reader = PdfReader(str(path_obj))
            return "".join(p.extract_text() or "" for p in reader.pages)
        except Exception as e:
            console.print(f"[red]Cannot read PDF content: {e}[/red]")
            raise typer.Exit(code=1)

    elif ext == ".docx":
        try:
            doc = docx.Document(str(path_obj))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            console.print(f"[red]Cannot read DOCX content: {e}[/red]")
            raise typer.Exit(code=1)

    elif ext == ".pptx":
        try:
            return _extract_pptx_text(str(path_obj))
        except Exception as e:
            console.print(f"[red]Cannot read PPTX content: {e}[/red]")
            raise typer.Exit(code=1)

    else:
        console.print(f"[red]Unsupported file type: {ext}[/red]")
        raise typer.Exit(code=1)
